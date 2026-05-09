# =========================================================
# FJIFE 2026 — Créateur de Présentation PowerPoint
# Rassemble les attestations générées en un seul fichier PPT
# prêt à être imprimé (1 attestation par slide)
# =========================================================

from pptx import Presentation
from pptx.util import Inches, Emu
from pptx.enum.shapes import MSO_SHAPE
import os
import glob

# =========================================================
# CONFIGURATION
# =========================================================

OUTPUT_FOLDER = "output"
PPT_FILENAME = "FJIFE26_Attestations_Ready_to_Print.pptx"

# Dimensions de slide en pouces (ratio 4:3 pour correspondre au template 4000x3000)
# A4 paysage ≈ 13.333 x 10 pouces (ratio 1.333 ≈ 4:3)
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(10)

# =========================================================
# CRÉATION DU POWERPOINT
# =========================================================

prs = Presentation()
prs.slide_width = SLIDE_WIDTH
prs.slide_height = SLIDE_HEIGHT

# Récupérer toutes les images PNG du dossier output, triées par nom
image_paths = sorted(glob.glob(os.path.join(OUTPUT_FOLDER, "*.png")))

if not image_paths:
    print(f"Aucune image trouvée dans '{OUTPUT_FOLDER}/'")
    print("Veuillez d'abord générer les attestations avec : python3 attestation_generator_script.py")
    exit(1)

print(f"Nombre d'attestations trouvées : {len(image_paths)}")

for img_path in image_paths:
    # Ajouter une slide vierge (layout vide)
    blank_layout = prs.slide_layouts[6]  # 6 = blank layout
    slide = prs.slides.add_slide(blank_layout)

    # Ajouter l'image en pleine page sans marges
    slide.shapes.add_picture(
        img_path,
        left=Inches(0),
        top=Inches(0),
        width=SLIDE_WIDTH,
        height=SLIDE_HEIGHT
    )

    print(f"  Slide ajoutée : {os.path.basename(img_path)}")

# Sauvegarder le fichier PPT
prs.save(PPT_FILENAME)

print(f"\nPrésentation créée avec succès : {PPT_FILENAME}")
print(f"Nombre total de slides : {len(image_paths)}")
print("Format : pleine page, prêt à imprimer.")
